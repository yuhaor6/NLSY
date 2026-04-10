#!/usr/bin/env python3
"""
Build lab_meeting_presentation.pptx
NLSY79 Labor Signaling Pipeline — Lab Meeting, March 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── Color palette ─────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x1B, 0x3A, 0x6B)
NAVY_MID  = RGBColor(0x2B, 0x5D, 0xB8)
BODY      = RGBColor(0x1A, 0x1A, 0x1A)
MUTED     = RGBColor(0x55, 0x55, 0x55)
FAINT     = RGBColor(0x88, 0x88, 0x88)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT  = RGBColor(0xF5, 0xF5, 0xF3)
CB_BG     = RGBColor(0xED, 0xF4, 0xFF)
CB_BR     = RGBColor(0x2B, 0x5D, 0xB8)
CG_BG     = RGBColor(0xEA, 0xF6, 0xEE)
CG_BR     = RGBColor(0x27, 0x7A, 0x4B)
CA_BG     = RGBColor(0xFF, 0xF8, 0xEA)
CA_BR     = RGBColor(0xA0, 0x70, 0x10)
CR_BG     = RGBColor(0xFF, 0xF0, 0xF0)
CR_BR     = RGBColor(0xB8, 0x30, 0x30)
DARK_NAVY = RGBColor(0x12, 0x28, 0x4E)
GREEN_DRK = RGBColor(0x1B, 0x5E, 0x3A)
RULE_CLR  = RGBColor(0xD4, 0xD4, 0xD4)
ALT_ROW   = RGBColor(0xF8, 0xF8, 0xF6)

# ─── Slide dimensions (13.333" × 7.5", standard widescreen) ───────────────────
W  = Inches(13.333)
H  = Inches(7.5)

# Layout constants
MX  = Inches(0.72)          # horizontal margin
CW  = W - 2 * MX            # content width ≈ 11.89"
GAP = Inches(0.33)           # column gap

# Body Y-start
BY_SUB = Inches(1.54)        # with subtitle tag
BY_NOS = Inches(1.33)        # without subtitle tag

# Body height (to bottom margin 0.45")
BH_SUB = H - BY_SUB - Inches(0.45)
BH_NOS = H - BY_NOS - Inches(0.45)

# Two-column widths
COL2 = (CW - GAP) / 2           # ≈ 5.78"
COL2_L = MX                      # left col x
COL2_R = MX + COL2 + GAP        # right col x

# 55/45 split
COL55 = (CW - GAP) * 0.55       # ≈ 6.34"
COL45 = (CW - GAP) * 0.45       # ≈ 5.19"
COL55_R = MX + COL55 + GAP      # right col x for 55/45

# Three-column widths
COL3 = (CW - 2 * GAP) / 3       # ≈ 3.74"
COL3_1 = MX
COL3_2 = MX + COL3 + GAP
COL3_3 = MX + 2 * (COL3 + GAP)

# ─── Helpers ───────────────────────────────────────────────────────────────────

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _rect(slide, left, top, width, height, fill, line=None):
    sh = slide.shapes.add_shape(1, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
    else:
        sh.line.fill.background()
    return sh

def _txb(slide, text, left, top, width, height,
         size=18, bold=False, italic=False, color=BODY,
         align=PP_ALIGN.LEFT, font='Calibri', wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.text_frame.word_wrap = wrap
    p = txb.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return txb

def _para(tf, text='', size=17, bold=False, italic=False, color=BODY,
          align=PP_ALIGN.LEFT, space_before=0, space_after=0, add=True):
    """Add a paragraph to a text frame. Returns (paragraph, run)."""
    if add:
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if text:
        r = p.add_run()
        r.text = text
        r.font.name = 'Calibri'
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        return p, r
    return p, None

def _mtxb(slide, left, top, width, height, wrap=True):
    """Create an empty multi-paragraph textbox."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.text_frame.word_wrap = wrap
    return txb

# ─── Slide-level components ────────────────────────────────────────────────────

def slide_header(slide, title, subtitle_tag=None):
    """Add standard header: (subtitle_tag), title, rule. Returns body_y."""
    if subtitle_tag:
        _txb(slide, subtitle_tag, MX, Inches(0.65), CW, Inches(0.28),
             size=12, bold=True, color=NAVY_MID)
        ty = Inches(0.90)
    else:
        ty = Inches(0.65)

    _txb(slide, title, MX, ty, CW - Inches(1.0), Inches(0.68),
         size=27, bold=True, color=NAVY)

    ry = ty + Inches(0.70)
    rule = slide.shapes.add_shape(1, MX, ry, CW, Pt(2))
    rule.fill.solid(); rule.fill.fore_color.rgb = NAVY_MID
    rule.line.fill.background()

    return ry + Inches(0.12)   # body starts here

def slide_num(slide, n, total=14):
    _txb(slide, f"{n} / {total}",
         W - Inches(1.2), H - Inches(0.35), Inches(1.1), Inches(0.25),
         size=12, color=FAINT, align=PP_ALIGN.RIGHT)

def col_label(slide, text, left, top, width):
    """Draw section label with bottom rule. Returns new top after label."""
    _txb(slide, text.upper(), left, top, width, Inches(0.26),
         size=11, bold=True, color=NAVY)
    rule = slide.shapes.add_shape(1, left, top + Inches(0.24), width, Pt(1.5))
    rule.fill.solid(); rule.fill.fore_color.rgb = RULE_CLR
    rule.line.fill.background()
    return top + Inches(0.38)

def callout(slide, text, left, top, width, height,
            bg=CB_BG, border=CB_BR, label=None, fs=15):
    """Left-bordered callout box."""
    bw = Inches(0.07)
    _rect(slide, left, top, bw, height, border)
    _rect(slide, left + bw, top, width - bw, height, bg)

    tx, ty = left + bw + Inches(0.12), top + Inches(0.10)
    tw, th = width - bw - Inches(0.20), height - Inches(0.18)
    txb = _mtxb(slide, tx, ty, tw, th)
    tf = txb.text_frame
    tf.word_wrap = True

    first = True
    if label:
        _, _ = _para(tf, label.upper(), size=10, bold=True, color=MUTED,
                     space_after=2, add=(not first))
        first = False

    _para(tf, text, size=fs, color=BODY, add=(not first))

def eq_box(slide, text, left, top, width):
    """Equation block with gray background. Returns bottom Y."""
    h = Inches(0.60)
    _rect(slide, left, top, width, h, BG_LIGHT)
    b = slide.shapes.add_shape(1, left, top, width, h)
    b.fill.background(); b.line.color.rgb = RULE_CLR; b.line.width = Pt(0.75)
    _txb(slide, text, left + Inches(0.2), top + Inches(0.08),
         width - Inches(0.4), h - Inches(0.14),
         size=15, color=NAVY, align=PP_ALIGN.CENTER, font='Cambria')
    return top + h

def result_box(slide, num, label, detail,
               left, top, width, height,
               bg=CB_BG, border=CB_BR, num_color=None):
    if num_color is None:
        num_color = NAVY
    bw = Inches(0.07)
    _rect(slide, left, top, bw, height, border)
    _rect(slide, left + bw, top, width - bw, height, bg)

    tx, ty = left + bw + Inches(0.10), top + Inches(0.08)
    tw, th = width - bw - Inches(0.18), height - Inches(0.14)
    txb = _mtxb(slide, tx, ty, tw, th)
    tf = txb.text_frame
    tf.word_wrap = True

    _para(tf, label.upper(), size=10, bold=True, color=MUTED, space_after=1, add=False)
    _para(tf, num, size=24, bold=True, color=num_color, space_after=1)
    _para(tf, detail, size=12, color=MUTED)

def bullets(slide, items, left, top, width, height, fs=16):
    """Multi-paragraph bullet list.
    Each item: str  → plain bullet
               (em_text, rest) → bold lead + normal text
               ('SUB', text)  → sub-bullet (indented, smaller)
    """
    txb = _mtxb(slide, left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True

    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(3)

        if isinstance(item, str):
            r = p.add_run(); r.text = '– '; r.font.name = 'Calibri'
            r.font.size = Pt(fs); r.font.bold = True; r.font.color.rgb = NAVY_MID
            r2 = p.add_run(); r2.text = item; r2.font.name = 'Calibri'
            r2.font.size = Pt(fs); r2.font.color.rgb = BODY
        elif isinstance(item, tuple) and item[0] == 'SUB':
            _, txt = item
            r = p.add_run(); r.text = '    · '; r.font.name = 'Calibri'
            r.font.size = Pt(fs - 1); r.font.color.rgb = FAINT
            r2 = p.add_run(); r2.text = txt; r2.font.name = 'Calibri'
            r2.font.size = Pt(fs - 1); r2.font.color.rgb = MUTED
        else:
            em, rest = item
            r = p.add_run(); r.text = '– '; r.font.name = 'Calibri'
            r.font.size = Pt(fs); r.font.bold = True; r.font.color.rgb = NAVY_MID
            r2 = p.add_run(); r2.text = em; r2.font.name = 'Calibri'
            r2.font.size = Pt(fs); r2.font.bold = True; r2.font.color.rgb = NAVY
            if rest:
                r3 = p.add_run(); r3.text = rest; r3.font.name = 'Calibri'
                r3.font.size = Pt(fs); r3.font.color.rgb = BODY

def tbl(slide, headers, rows, left, top, width, height,
        hl_last=False, col_widths=None, row_fs=13):
    nc = len(headers)
    nr = len(rows) + 1
    t = slide.shapes.add_table(nr, nc, left, top, width, height).table
    if col_widths:
        for i, cw in enumerate(col_widths): t.columns[i].width = cw

    # Header
    for j, h in enumerate(headers):
        cell = t.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT if j > 0 else PP_ALIGN.LEFT
        run = p.add_run(); run.text = h
        run.font.bold = True; run.font.size = Pt(12)
        run.font.color.rgb = WHITE; run.font.name = 'Calibri'

    # Rows
    for i, row in enumerate(rows):
        is_hl = hl_last and (i == len(rows) - 1)
        for j, val in enumerate(row):
            cell = t.cell(i + 1, j)
            if is_hl:
                cell.fill.solid(); cell.fill.fore_color.rgb = CB_BG
            elif i % 2 == 1:
                cell.fill.solid(); cell.fill.fore_color.rgb = ALT_ROW
            else:
                cell.fill.background()
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if j > 0 else PP_ALIGN.LEFT
            run = p.add_run(); run.text = str(val)
            run.font.size = Pt(row_fs)
            run.font.bold = is_hl
            run.font.color.rgb = NAVY if is_hl else BODY
            run.font.name = 'Calibri'
    return t

# ──────────────────────────────────────────────────────────────────────────────
# BUILD PRESENTATION
# ──────────────────────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── SLIDE 1: Title ─────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg = s.background.fill; bg.solid(); bg.fore_color.rgb = DARK_NAVY

# Top label bar
_txb(s, 'Internal Lab Meeting', MX, Inches(0.55), Inches(3.5), Inches(0.28),
     size=11, bold=True, color=RGBColor(0xFF,0xFF,0xFF), font='Calibri')
_txb(s, 'Sztutman Labor & Tax Lab  ·  CMU  ·  March 2026',
     W - MX - Inches(5), Inches(0.55), Inches(5), Inches(0.28),
     size=11, bold=True, color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.RIGHT, font='Calibri')
rule = s.shapes.add_shape(1, MX, Inches(0.85), CW, Pt(1))
rule.fill.solid(); rule.fill.fore_color.rgb = RGBColor(0xFF,0xFF,0xFF)
rule.line.fill.background()

# Main title
txb = _mtxb(s, MX, Inches(1.8), CW, Inches(1.8))
tf = txb.text_frame; tf.word_wrap = True
_, _ = _para(tf,
    'Wage Growth, Skill, and Signaling:\nBuilding an NLSY79 Empirical Pipeline',
    size=46, bold=True, color=WHITE, font='Calibri', add=False)

# Subtitle
_txb(s,
    'Pipeline progress report — data construction, ETI estimation,\n'
    'and wage elasticity of income via BSX-adapted IV',
    MX, Inches(3.75), CW, Inches(0.8),
    size=20, italic=True, color=RGBColor(0xBB,0xCC,0xEE), font='Calibri')

# Bottom rule
rule2 = s.shapes.add_shape(1, MX, Inches(4.75), CW, Pt(1))
rule2.fill.solid(); rule2.fill.fore_color.rgb = RGBColor(0x55,0x77,0xAA)
rule2.line.fill.background()

# Presenter info
_txb(s, 'Yuhao Chen', MX, Inches(4.95), CW * 0.6, Inches(0.38),
     size=20, bold=True, color=RGBColor(0xEE,0xEE,0xFF), font='Calibri')
_txb(s, 'Carnegie Mellon University  ·  Advisor: Professor André Sztutman',
     MX, Inches(5.35), CW * 0.75, Inches(0.30),
     size=15, color=RGBColor(0x88,0x99,0xBB), font='Calibri')
_txb(s, 'March 2026  ·  NLSY79 Empirical Pipeline, Stage 1',
     MX, Inches(5.65), CW * 0.75, Inches(0.30),
     size=15, color=RGBColor(0x88,0x99,0xBB), font='Calibri')

# Pill
_rect(s, W - MX - Inches(2.3), Inches(4.95), Inches(2.3), Inches(0.75),
      RGBColor(0x25,0x3A,0x62))
_txb(s, 'NLSY79 First\nNLSY97 to Follow',
     W - MX - Inches(2.2), Inches(5.02), Inches(2.1), Inches(0.62),
     size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font='Calibri')

slide_num(s, 1)

# ── SLIDE 2: Research Question ─────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
by = slide_header(s, 'Research Question and Why It Matters')

# Left col (55%)
x1, w1 = COL2_L, COL55
cy = col_label(s, 'The Question', x1, by, w1)
_txb(s, 'Early-career wage growth reflects what, exactly?',
     x1, cy, w1, Inches(0.35), size=17, italic=True, color=BODY)
cy += Inches(0.40)
bullets(s, [
    ('Skill accumulation — ', 'workers become more productive over time through learning by doing and training'),
    ('Signaling — ', 'employers update beliefs about workers\' pre-existing ability; wages grow as uncertainty resolves'),
    ('Some combination — ', 'the two channels coexist and are difficult to separate empirically'),
], x1, cy, w1, Inches(1.30), fs=17)
cy += Inches(1.40)
callout(s,
    'If signaling is the dominant mechanism, labor taxes must correct for more than standard '
    'labor-supply distortions. Sztutman (2024) shows the optimal tax wedge depends on the '
    'relative weight of signaling versus productive effort — a wedge not captured by ETI alone.',
    x1, cy, w1, Inches(1.30), bg=CB_BG, border=CB_BR,
    label='Why the distinction matters for taxation', fs=14)

# Right col (45%)
x2, w2 = COL55_R, COL45
cy2 = col_label(s, 'This Stage: Build the Pipeline', x2, by, w2)
bullets(s, [
    'This is a first-stage empirical build, not a final welfare calculation',
    'Primary objective: construct a coherent NLSY79 person-year panel and evaluate which empirical branches are feasible',
    'Three branches under development:',
    ('SUB', 'Gruber–Saez-style ETI estimation'),
    ('SUB', 'BSX-adapted wage elasticity of income'),
    ('SUB', 'Exploratory skill vs. signaling diagnostics'),
    'State tax variation not yet available — confidential geography pending',
    'NLSY79 first; NLSY97 to follow',
], x2, cy2, w2, Inches(2.80), fs=15)
cy2 += Inches(3.00)
callout(s,
    'Today\'s goal: present what has been built, report what the pipeline produces, '
    'and ask for feedback on priorities.',
    x2, cy2, w2, Inches(0.65), bg=CA_BG, border=CA_BR, fs=15)

slide_num(s, 2)

# ── SLIDE 3: Gruber–Saez Lit Map ───────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
by = slide_header(s, 'Gruber & Saez (2002) — What I Borrowed and Why',
                  subtitle_tag='Literature Map, Part 1')

# Left col
x1, w1, cy1 = COL2_L, COL2, by
cy1 = col_label(s, 'What They Did', x1, cy1, w1)
bullets(s, [
    'Estimate the elasticity of taxable income (ETI) with respect to the net-of-tax rate using NBER TAXSIM and U.S. panel data',
    'Use tax reforms (ERTA 1981, TRA 1986) as natural experiments generating variation in marginal rates',
    ('Key identification move: ', 'simulated instrument — predict each person\'s tax change holding pre-reform income fixed'),
    'Three-year first differences to allow behavioral adjustment',
    'Real income floor ($10K, 1984 dollars) to exclude workers with very unstable income',
    'Marital stability requirement across the period',
], x1, cy1, w1, Inches(3.80), fs=16)

# Right col
x2, w2, cy2 = COL2_R, COL2, by
cy2 = col_label(s, 'What I Borrowed and Why', x2, cy2, w2)
bullets(s, [
    ('Simulated instrument logic: ', 'TAXSIM35 computes the counterfactual MTR holding base-year income fixed → isolates policy variation from behavioral response'),
    ('9-piece linear spline on log(z_t), ', 'knots at deciles 10–90 — income growth is nonlinear in starting income; without this, mean reversion among low earners is conflated with tax effects'),
    ('Three-year differences ', '(annual); two-year differences (biennial) — matches survey cadence; 3-year lags would skip biennial waves'),
    ('Clustering by individual ', 'for standard errors'),
], x2, cy2, w2, Inches(2.60), fs=16)
cy2 += Inches(2.75)
callout(s,
    'No state FIPS codes in current extract → state MTR = 0. Results are federal-only. '
    'State variation pending confidential geography.',
    x2, cy2, w2, Inches(0.80), bg=CA_BG, border=CA_BR,
    label='Current limitation', fs=14)

slide_num(s, 3)

# ── SLIDE 4: BSX + Sztutman Lit Map ────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
by = slide_header(s, 'Becko, Sztutman & Xia (2024) — What I Borrowed, and What Is Deferred',
                  subtitle_tag='Literature Map, Part 2')

# Left col
x1, w1, cy1 = COL2_L, COL2, by
cy1 = col_label(s, 'BSX (2024): The Instrument and the Target', x1, cy1, w1)
bullets(s, [
    ('BSX estimate the wage elasticity of income (ε_w) — ', 'how broad taxable income responds to an exogenous change in hourly wages'),
    ('Identification: leave-one-out (LOO) industry-cohort mean wage change — ', 'the average wage growth in a worker\'s industry excluding that worker isolates demand-side shocks from individual choices'),
    'Target equation (BSX Eq. 30): Δlog(z) = α + ε_w · Δlog(w) + year FE + marital status + splines + u',
], x1, cy1, w1, Inches(1.90), fs=16)
cy1 += Inches(2.05)
callout(s,
    'BSX instrument = CHANGE in LOO mean wage from t to t+3, not the LEVEL at t+3. '
    'Using the level conflates persistent industry wage premia (selection channel) with '
    'transitory demand shocks (identification channel). The change instrument strips out '
    'the level-selection bias. Corrected estimate dropped from 1.26 → 0.98.',
    x1, cy1, w1, Inches(1.50), bg=CB_BG, border=CB_BR,
    label='The correction I made (BSX Eq. 31)', fs=14)

# Right col
x2, w2, cy2 = COL2_R, COL2, by
cy2 = col_label(s, 'Sztutman (2024, JMP): The Broader Framework', x2, cy2, w2)
bullets(s, [
    'Structural motivation: signaling + taxation → labor wedge depending on the ratio of signaling to productive effort',
    ('Welfare formula requires ε_w and a productivity elasticity η^P — ', 'η^P distinguishes genuine skill formation from revelation of pre-existing ability'),
], x2, cy2, w2, Inches(1.20), fs=16)
cy2 += Inches(1.35)
cy2 = col_label(s, 'What Is Intentionally Deferred', x2, cy2, w2)
bullets(s, [
    'Computing the tax wedge χ = 1 + ε_w/η^P — requires credible identification of η^P',
    'Full structural calibration of the signaling model',
    'Skill vs. signaling decomposition — only exploratory diagnostics exist now',
    'Any welfare or optimal-tax claims',
], x2, cy2, w2, Inches(1.40), fs=16)
cy2 += Inches(1.55)
callout(s,
    'This stage: Build the empirical engine. Evaluate which branches can credibly '
    'deliver the inputs the theory needs.',
    x2, cy2, w2, Inches(0.65), bg=CA_BG, border=CA_BR, fs=15)

slide_num(s, 4)

# ── SLIDE 5: Data Built ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
by = slide_header(s, 'What the Pipeline Built: NLSY79 Person-Year Panel')

# Col 1: Core Panel + Outcomes
x1, w1, cy1 = COL3_1, COL3, by
cy1 = col_label(s, 'Core Panel', x1, cy1, w1)
bullets(s, [
    '12,686 respondents',
    '46 survey waves (1979–2020)',
    '583,556 person-year obs',
    'Annual: 1978–1993 | Biennial: 1995–2019+',
    'Long format: xtset taxsimid year',
], x1, cy1, w1, Inches(1.40), fs=15)
cy1 += Inches(1.52)
cy1 = col_label(s, 'Key Outcomes', x1, cy1, w1)
bullets(s, [
    'pwages, swages — primary & spouse wages',
    'psemp, ssemp — self-employment income',
    'pui, sui, gssi, pensions',
    'broad_income = sum of all components',
], x1, cy1, w1, Inches(1.20), fs=15)
cy1 += Inches(1.32)
cy1 = col_label(s, 'Ability & Background', x1, cy1, w1)
bullets(s, [
    'afqt_pct — 1980/1989/2006 revisions',
    'hgc — year-specific HGC (corrected)',
    'occ_broad, ind_broad — 1979–1993 only',
], x1, cy1, w1, Inches(0.90), fs=15)

# Col 2: Experience + TAXSIM
x2, w2, cy2 = COL3_2, COL3, by
cy2 = col_label(s, 'Experience & Hours', x2, cy2, w2)
bullets(s, [
    'pot_exp — potential experience (corrected)',
    'cumhrs — cumulative hours (interpolated)',
    'hrs — annual hours, all jobs',
], x2, cy2, w2, Inches(0.90), fs=15)
cy2 += Inches(1.02)
cy2 = col_label(s, 'TAXSIM-Ready Variables', x2, cy2, w2)
bullets(s, [
    'mstat (NLSY → TAXSIM mapping, corrected)',
    'page, sage — income-year ages',
    'depx — number of dependents',
    'Federal MTR & tax liability via TAXSIM35',
    'State MTR: pending (no state FIPS yet)',
], x2, cy2, w2, Inches(1.50), fs=15)
cy2 += Inches(1.62)
callout(s,
    'Annual surveys 1979–1994 → income years 1978–1993. Biennial surveys 1996, 1998, … '
    '→ income years 1995, 1997, …. Education and age variables realigned to income year.',
    x2, cy2, w2, Inches(1.10), bg=CB_BG, border=CB_BR,
    label='Annual vs. Biennial Structure', fs=14)

# Col 3: Construction Fixes
x3, w3, cy3 = COL3_3, COL3, by
cy3 = col_label(s, 'Critical Construction Fixes', x3, cy3, w3)
bullets(s, [
    ('Fix 1 — Potential experience: ', 'uses current-year hgc; prevents experience inflation for students'),
    ('Fix 2 — Education in long format: ', 'year-specific hgc created; biennial waves realigned to income year'),
    ('Fix 3 — Cumulative hours: ', 'linear interpolation for biennial gap years; plateau artifact corrected'),
    ('Fix 4 — Age alignment: ', 'page = income-year age (survey age − 1 for biennial periods)'),
    ('Fix 5 — Marital status: ', 'corrected NLSY → TAXSIM mapping (1=Single, 2=MFJ)'),
], x3, cy3, w3, Inches(3.80), fs=15)

slide_num(s, 5)

# ── SLIDE 6: ETI Method ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
by = slide_header(s, 'Tax Response Analysis: Method and Key Implementation Choices',
                  subtitle_tag='Section A — Adapted from Gruber & Saez (2002)')

ey = eq_box(s,
    'Δlog(z_{i,t+k}) = α + ε · Δlog(1 − τ_{i,t+k}) + δ_t + γ·M_{i,t} + Σ_j θ_j·S_j(log z_{i,t}) + u_{i,t}',
    MX, by, CW)
by = ey + Inches(0.15)

# Left col
x1, w1, cy1 = COL2_L, COL2, by
cy1 = col_label(s, 'Instrument & Period Structure', x1, cy1, w1)
bullets(s, [
    ('Simulated instrument: ', 'TAXSIM35 computes MTR at base-year income under both-year tax law, holding income fixed — isolates policy variation from behavioral response'),
    ('Period pairing: ', 'k = 3 years (annual, 1978–1993); k = 2 years (biennial, 1995–2019). Two-year lag for biennial matches survey frequency — a 3-year lag would bridge a non-survey year'),
    ('Broad income: ', 'wages + spouse wages + self-employment + UI + SSDI + pensions; real 1984 dollars (CPI-U)'),
    ('Restrictions: ', 'real income floor ≥ $10K (1984 $); marital stability; EITC phase-out excluded'),
], x1, cy1, w1, Inches(3.00), fs=16)

# Right col
x2, w2, cy2 = COL2_R, COL2, by
cy2 = col_label(s, 'Key Implementation Choices', x2, cy2, w2)
bullets(s, [
    ('9-piece linear spline on log(z_t), ', 'knots at deciles 10–90 — separates tax variation from income dynamics; mean reversion among low earners otherwise conflated with behavioral response'),
    ('Standard errors ', 'clustered by individual (taxsimid)'),
    ('Year fixed effects (δ_t) ', 'absorb aggregate macroeconomic income trends'),
], x2, cy2, w2, Inches(2.00), fs=16)
cy2 += Inches(2.15)
callout(s,
    'No state FIPS in current data extract → state MTR = 0 throughout. Results reflect '
    'federal tax variation only. State component will be added once confidential geography '
    'arrives — the single most important missing input for this branch.',
    x2, cy2, w2, Inches(1.10), bg=CA_BG, border=CA_BR,
    label='Current Limitation — Federal Only', fs=14)

slide_num(s, 6)

# ── SLIDE 7: ETI Results ────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
by = slide_header(s, 'ETI Estimates: Pipeline Running, Sign Instability to Diagnose',
                  subtitle_tag='Section A — Results')

# Left col: table
x1, w1 = COL2_L, COL55
tbl(s,
    ['', 'Annual (1978–1993)', 'Biennial (1995–2019)'],
    [
        ['Period length', '3-year lag', '2-year lag'],
        ['Final sample', '31,832', '47,201'],
        ['Age at baseline', '~17–35', '~31–62'],
        ['Tax reforms', 'ERTA 1981, TRA 1986', 'EGTRRA 2001, JGTRRA 2003, TCJA 2017'],
        ['First-stage F', '1,349', '988'],
        ['ETI estimate', '−0.173', '+0.544'],
    ],
    x1, by, w1, Inches(2.30),
    hl_last=True,
    col_widths=[Inches(2.0), Inches(2.0), Inches(2.34)],
    row_fs=14)
by_img = by + Inches(2.45)
_txb(s,
    '[First-stage scatter plots: output/Fig1a_FirstStage_Annual.png, Fig1b_FirstStage_Biennial.png]',
    x1, by_img, w1, Inches(0.5),
    size=12, italic=True, color=FAINT)

# Right col
x2, w2, cy2 = COL55_R, COL45, by
cy2 = col_label(s, 'Current Interpretation', x2, cy2, w2)
_txb(s,
    'Both periods have very strong first stages (F ≫ 10), confirming the simulated '
    'instrument is highly predictive. The pipeline runs end-to-end. However, the '
    'sign reversal between periods is the central open question.',
    x2, cy2, w2, Inches(0.95),
    size=15, color=BODY)
cy2 += Inches(1.10)
callout(s,
    '– Young workers (17–35) have near-zero or negative labor supply elasticity — '
    'life-cycle transitions dominate\n'
    '– Prime-age workers (31–62) have positive elasticity consistent with literature\n'
    '– Different tax reform regimes; different income distributions\n'
    '– Federal-only variation may amplify noise in the annual period',
    x2, cy2, w2, Inches(1.50), bg=CA_BG, border=CA_BR,
    label='Why the signs differ — leading hypotheses', fs=14)
cy2 += Inches(1.65)
callout(s,
    'The annual ETI of −0.173 should not be interpreted as a causal behavioral elasticity. '
    'The sign instability needs diagnosis through subgroup analysis, income-group breakdowns, '
    'and reform-by-reform comparison before any interpretation is possible.',
    x2, cy2, w2, Inches(1.30), bg=CR_BG, border=CR_BR,
    label='Bottom line: not a final estimate', fs=14)

slide_num(s, 7)

# ── SLIDE 8: BSX Method ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
by = slide_header(s, 'Wage Elasticity of Income: Method and the Instrument Correction',
                  subtitle_tag='Section B — Adapted from Becko, Sztutman & Xia (2024)')

ey = eq_box(s,
    'Δlog(z_{i,t+3}) = α + ε_w · Δlog(w_{i,t+3}) + δ_t + γ·M_{i,t} + Σ_j θ_j·S_j(log z_{i,t}) + u_{i,t}',
    MX, by, CW)
by = ey + Inches(0.15)

# Left col (55%)
x1, w1, cy1 = COL2_L, COL55, by
cy1 = col_label(s, 'Instrument (BSX Eq. 31, Corrected)', x1, cy1, w1)
bullets(s, [
    ('Leave-one-out (LOO) industry cohort wage change: ', 'average log-wage change in worker i\'s industry from t to t+3, excluding worker i — isolates industry-level demand shocks from individual choices'),
    'Validity threshold: LOO cohort N ≥ 5 per industry × year cell',
    'Sample restricted to 1979–1993 (ind_broad available only in this period)',
], x1, cy1, w1, Inches(1.70), fs=16)
cy1 += Inches(1.85)
callout(s,
    'Using the LOO wage LEVEL at t+3 conflates persistent industry wage premia '
    '(selection channel) with transitory demand shocks (identification channel). '
    'Workers sorted into high-wage industries earn more through promotions — unrelated '
    'to the wage shock. The CHANGE from t to t+3 strips out this level-selection bias. '
    'Corrected estimate: ε_w = 0.98 vs. 1.26 (level instrument).',
    x1, cy1, w1, Inches(1.30), bg=CB_BG, border=CB_BR,
    label='Why change, not level (the key correction)', fs=14)

# Right col (45%)
x2, w2, cy2 = COL55_R, COL45, by
cy2 = col_label(s, 'Wage Measure & Income Target', x2, cy2, w2)
bullets(s, [
    ('Hourly wage: ', 'w = pwages / hrs (primary wages / total hours, all jobs). For multi-job workers this underestimates the primary-job rate — introduces measurement error in the endogenous regressor'),
    'Attenuation is modest: first-stage F = 54.87 despite the noise. BSX use directly surveyed hourly rates (SIPP)',
    ('DV: ', 'Δlog(broad_income); robustness uses Δlog(pwages) alone — both reported'),
], x2, cy2, w2, Inches(2.10), fs=15)
cy2 += Inches(2.25)
callout(s,
    'Mean = 0.205, SD = 0.071, range [0.03, 0.40]. Variation is continuous and '
    'well-behaved across industries and years.',
    x2, cy2, w2, Inches(0.85), bg=CA_BG, border=CA_BR,
    label='Instrument distribution', fs=14)

slide_num(s, 8)

# ── SLIDE 9: BSX Results ────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
by = slide_header(s, 'Wage Elasticity of Income: ε_w ≈ 0.98 (Corrected IV)',
                  subtitle_tag='Section B — Results')

# Left col: result boxes in a 2×3 grid
x1, w1 = COL2_L, COL55
rb_w = (w1 - GAP) / 2
rb_h = Inches(0.90)

result_box(s, '29,379', 'Final sample', 'obs, 1979–1993',
           x1, by, rb_w, rb_h, bg=CG_BG, border=CG_BR, num_color=GREEN_DRK)
result_box(s, '54.87', 'First-stage F', 'β = 0.634, SE = 0.086',
           x1 + rb_w + GAP, by, rb_w, rb_h, bg=CG_BG, border=CG_BR, num_color=GREEN_DRK)

ry2 = by + rb_h + Inches(0.10)
result_box(s, '0.710', 'OLS ε_w', 'SE = 0.026  (attenuated)',
           x1, ry2, rb_w, rb_h)
result_box(s, '0.981', '★ 2SLS ε_w  (main result)', 'SE = 0.121   N = 29,379',
           x1 + rb_w + GAP, ry2, rb_w, rb_h,
           bg=RGBColor(0xE6,0xED,0xF9), border=NAVY, num_color=NAVY)

ry3 = ry2 + rb_h + Inches(0.10)
result_box(s, '0.981', 'LIML ε_w', '= 2SLS → no weak-IV bias',
           x1, ry3, rb_w, rb_h)
result_box(s, '1.105', 'Robustness (Δlog pwages DV)', 'SE = 0.077',
           x1 + rb_w + GAP, ry3, rb_w, rb_h, bg=CA_BG, border=CA_BR, num_color=NAVY)

# Right col
x2, w2, cy2 = COL55_R, COL45, by
cy2 = col_label(s, 'Interpretation', x2, cy2, w2)
_txb(s,
    'A 1% exogenous increase in hourly wages leads to approximately a 1% increase in '
    'broad taxable income — consistent with the theoretical prior of inelastic hours '
    'adjustment (implied hours elasticity ≈ −0.02). OLS = 0.71 vs. 2SLS = 0.98: '
    'attenuation from measurement error in the constructed hourly wage (pwages/hrs) '
    'is consistent with this gap.',
    x2, cy2, w2, Inches(1.30), size=15, color=BODY)
cy2 += Inches(1.45)
callout(s,
    'Level instrument (old): ε_w = 1.26 (SE 0.157)\n'
    'Change instrument (corrected): ε_w = 0.98 (SE 0.121)\n'
    'The downward correction is consistent with the level instrument having captured '
    'persistent industry wage premia rather than transitory demand shocks.',
    x2, cy2, w2, Inches(1.30), bg=CG_BG, border=CG_BR,
    label='The instrument correction mattered', fs=14)
cy2 += Inches(1.45)
callout(s,
    'Still exploratory. Wage construction caveat (primary wages / total hours) means '
    'this is an approximation of BSX\'s directly surveyed hourly rates. Pending review '
    'of multi-job workers\' influence on results.',
    x2, cy2, w2, Inches(1.00), bg=CA_BG, border=CA_BR, fs=14)

slide_num(s, 9)

# ── SLIDE 10: Exploratory Diagnostics ──────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
by = slide_header(s, 'Descriptive Diagnostics: Industry Structure and Wage Profiles',
                  subtitle_tag='Section C — Exploratory  [Diagnostic, Not Causal]')

# Left col: cell structure
x1, w1, cy1 = COL2_L, COL2, by
cy1 = col_label(s, 'Occupation × Industry Cell Structure (1979–1993)', x1, cy1, w1)

rb_sw = (w1 - 2 * GAP) / 3
rb_sh = Inches(0.80)
result_box(s, '107', 'Total cells', 'occ × ind pairs',
           x1, cy1, rb_sw, rb_sh)
result_box(s, '86%', 'Cells ≥ 30 obs', '92 of 107 usable',
           x1 + rb_sw + GAP, cy1, rb_sw, rb_sh, bg=CG_BG, border=CG_BR, num_color=GREEN_DRK)
result_box(s, '317', 'Median cell', 'obs (mean: 978)',
           x1 + 2*(rb_sw + GAP), cy1, rb_sw, rb_sh)

cy1 += rb_sh + Inches(0.15)
bullets(s, [
    ('Largest cell: ', 'Operatives × Manufacturing — 8,594 obs, 3,043 individuals, mean wage $13,251'),
    ('Highest wages: ', 'Professional/Technical × Manufacturing — mean $28,279'),
    ('Lowest wages: ', 'Service Workers × Personal Services — mean $6,833'),
    '~5× wage differential across cells supports meaningful cross-industry variation in the LOO instrument',
], x1, cy1, w1, Inches(2.00), fs=16)

# Right col: ability patterns
x2, w2, cy2 = COL2_R, COL2, by
cy2 = col_label(s, 'Wage-Hours and Ability Patterns', x2, cy2, w2)
bullets(s, [
    ('Cumulative hours elasticity: ', '10% increase in cumhrs → ~9.4% higher wages (strong signal of skill-formation story)'),
    ('Recent work intensity: ', '10% increase → ~9.7% higher wages'),
    ('AFQT effect: ', '1 SD increase → ~27% higher wages — ability premium is large and persistent'),
    ('AFQT × experience: ', 'AFQT effect appears to increase with experience (supports human capital alongside selection)'),
], x2, cy2, w2, Inches(2.40), fs=16)
cy2 += Inches(2.55)
callout(s,
    'These patterns are purely descriptive and motivate the decomposition question, '
    'but they do not identify signaling vs. skill formation. They are hypothesis-generating '
    'inputs, not causal evidence.',
    x2, cy2, w2, Inches(0.95), bg=CA_BG, border=CA_BR, fs=15)

slide_num(s, 10)

# ── SLIDE 11: What Is Working ──────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
by = slide_header(s, 'What the Pipeline Has Achieved')

items = [
    ('NLSY79 Panel Construction Is Coherent',
     'The 583,556-observation person-year panel has year-specific education, corrected potential '
     'experience, interpolated cumulative hours, and income-year age alignment. Five construction '
     'fixes verified against TAXSIM output.'),
    ('ETI Pipeline Runs End-to-End',
     'Both annual (N = 31,832) and biennial (N = 47,201) branches produce estimates with strong '
     'first stages (F = 1,349 and F = 988). Sign difference between periods is a diagnostic '
     'question, not a pipeline failure — the machinery works.'),
    ('Corrected BSX Branch Produces Interpretable IV Results',
     'After correcting the instrument from a level to a change, the LOO wage elasticity estimate '
     'is ε_w = 0.98 (SE = 0.12, F = 54.87), LIML = 2SLS confirming no weak-instrument bias. '
     'The estimate is economically plausible.'),
    ('Exploratory Cell Analysis Confirms Feasible Granularity',
     '107 occupation × industry cells for 1979–1993, with 86% meeting a 30-obs minimum. '
     'Substantial wage heterogeneity (~5×) supports the LOO identification strategy.'),
]

item_h = Inches(1.15)
for i, (strong, detail) in enumerate(items):
    iy = by + i * (item_h + Inches(0.10))
    # Check circle
    circ = s.shapes.add_shape(9, MX, iy + Inches(0.08), Inches(0.38), Inches(0.38))
    circ.fill.solid(); circ.fill.fore_color.rgb = CG_BG
    circ.line.color.rgb = CG_BR; circ.line.width = Pt(1.5)
    _txb(s, '✓', MX, iy + Inches(0.06), Inches(0.38), Inches(0.38),
         size=14, bold=True, color=CG_BR, align=PP_ALIGN.CENTER)

    # Content
    tx = MX + Inches(0.52)
    tw = CW - Inches(0.52)
    txb = _mtxb(s, tx, iy, tw, item_h)
    tf = txb.text_frame; tf.word_wrap = True
    _para(tf, strong, size=18, bold=True, color=NAVY, add=False, space_after=2)
    _para(tf, detail, size=16, color=MUTED)

slide_num(s, 11)

# ── SLIDE 12: Questions for Lab ────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
by = slide_header(s, 'What I Am Unsure About — Questions for the Lab')

questions = [
    ('ETI vs. BSX branch: where to focus next?',
     'The ETI branch has an unresolved sign reversal between annual and biennial periods. The BSX '
     'branch has a cleaner current estimate (ε_w ≈ 0.98) but is limited to 1979–1993 (no ind_broad '
     'thereafter). Should I prioritize diagnosing ETI instability or refining BSX specification first?'),
    ('What is the right bar for the BSX NLSY79 adaptation?',
     'The wage measure (primary wages / total hours) is an imperfect proxy, and the sample is '
     'constrained to 1979–1993. At what point is this adaptation credible alongside BSX\'s '
     'SIPP-based results — or does a better wage measure need to be constructed first?'),
    ('How much investment in the skill/signaling branch now?',
     'The exploratory diagnostics are descriptive only. Full identification requires '
     'Altonji–Pierret-style regressions or occupation-switching tests. How much priority does '
     'this branch merit at this stage relative to ETI and BSX?'),
    ('Which diagnostics for the annual vs. biennial ETI instability?',
     'Leading candidates: (a) subgroup ETI by age, (b) income-tertile breakdown, '
     '(c) reform-by-reform comparison, (d) adding age fixed effects. '
     'Which is most informative for the research agenda?'),
]

item_h = Inches(1.10)
for i, (strong, detail) in enumerate(questions):
    iy = by + i * (item_h + Inches(0.08))
    # Number circle
    circ = s.shapes.add_shape(9, MX, iy + Inches(0.05), Inches(0.38), Inches(0.38))
    circ.fill.solid(); circ.fill.fore_color.rgb = NAVY
    circ.line.fill.background()
    _txb(s, str(i + 1), MX, iy + Inches(0.05), Inches(0.38), Inches(0.38),
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    tx = MX + Inches(0.52)
    tw = CW - Inches(0.52)
    txb = _mtxb(s, tx, iy, tw, item_h)
    tf = txb.text_frame; tf.word_wrap = True
    _para(tf, strong, size=18, bold=True, color=NAVY, add=False, space_after=2)
    _para(tf, detail, size=16, color=MUTED)

slide_num(s, 12)

# ── SLIDE 13: Appendix A — BSX Waterfall ───────────────────────────────────────
s = prs.slides.add_slide(BLANK)
by = slide_header(s, 'BSX Branch: Sample Restriction Waterfall',
                  subtitle_tag='Appendix A')
_txb(s, 'APPENDIX', W - MX - Inches(1.4), Inches(0.65), Inches(1.4), Inches(0.28),
     size=11, bold=True, color=MUTED, align=PP_ALIGN.RIGHT)

# Left: waterfall table
x1, w1 = COL2_L, COL55
tbl(s,
    ['Filter applied', 'N remaining', '% of prior'],
    [
        ['All 1979–1993 observations',        '190,290', '—'],
        ['With positive pwages',               '119,347', '62.7%'],
        ['With positive hours (all jobs)',      '116,070', '97.3%'],
        ['Hours quality filter (520–4,160/yr)', '103,661', '89.3%'],
        ['Valid log(hourly wage) AND ind_broad', '95,915', '92.5%'],
        ['Valid wages at t AND t+3',            '63,955', '66.7%'],
        ['Valid ind_broad at t',                '59,361', '92.8%'],
        ['Marital stability across period',     '49,165', '82.8%'],
        ['Income floor ≥ $10K (1984 $)',        '29,379', '59.8%'],
        ['Final analysis sample',              '29,379', '—'],
    ],
    x1, by, w1, Inches(3.60),
    hl_last=True,
    col_widths=[Inches(3.5), Inches(1.5), Inches(1.34)],
    row_fs=13)
_txb(s,
    'Note: LOO cohort ≥ 5 filter does not additionally trim (final step nonbinding).',
    x1, by + Inches(3.70), w1, Inches(0.28),
    size=12, italic=True, color=FAINT)

# Right col: commentary
x2, w2, cy2 = COL55_R, COL45, by
cy2 = col_label(s, 'Key Attrition Points', x2, cy2, w2)
bullets(s, [
    ('Largest drop: ', 'requiring valid wages at both t and t+3 removes 34% — reflects respondents who exit employment or stop being surveyed over the 3-year window'),
    ('Second largest: ', 'real income floor ($10K, 1984 $) retains only 59.8% of panel with valid wages — intentional, restricts to committed labor-force participants'),
    ('Hours quality filter ', '(10–80 hrs/wk) removes 10.7% — mostly very short spells and top-end data errors'),
    'Final N = 29,379 represents attached, prime-working-age workers — appropriate for ε_w estimation but limits external validity',
], x2, cy2, w2, Inches(2.50), fs=15)
cy2 += Inches(2.65)
callout(s,
    'Annual period: 202,976 → 31,832 final (15.7% retention). Biennial: → 47,201 final. '
    'The income floor accounts for most of the attrition in both branches.',
    x2, cy2, w2, Inches(0.90), bg=CA_BG, border=CA_BR,
    label='ETI Annual Branch (for comparison)', fs=14)

slide_num(s, 13)

# ── SLIDE 14: Appendix B — Variable Construction ───────────────────────────────
s = prs.slides.add_slide(BLANK)
by = slide_header(s, 'Variable Construction: Key Decisions and Rationale',
                  subtitle_tag='Appendix B')
_txb(s, 'APPENDIX', W - MX - Inches(1.4), Inches(0.65), Inches(1.4), Inches(0.28),
     size=11, bold=True, color=MUTED, align=PP_ALIGN.RIGHT)

tbl(s,
    ['Variable', 'Construction', 'Issue / Why it matters', 'What was done'],
    [
        ['pot_exp', 'age − hgc − 6',
         'If hgc = lifetime max, experience is inflated for students still in school',
         'Uses current-year hgc from each survey wave, not a fixed lifetime maximum'],
        ['hgc (long format)', 'Year-specific highest grade completed',
         'Biennial surveys report income from prior year; using survey-year education creates a one-year lag',
         'Biennial waves realigned: hgc_1996 → hgc_1995, hgc_1998 → hgc_1997, etc.'],
        ['cumhrs', 'Running sum of hrs across years',
         'Post-1994 biennial gaps leave even years unobserved, causing plateau artifacts',
         'Linear interpolation between adjacent survey years; plateau fix applied'],
        ['page', 'Primary taxpayer age',
         'TAXSIM expects age in the income year; biennial surveys happen in income+1',
         'Income-year age = survey age − 1 for biennial periods'],
        ['mstat', 'Marital status (TAXSIM)',
         'NLSY codes 1=Married / 2=Separated differ from TAXSIM 1=Single / 2=MFJ',
         'Corrected mapping: MFJ → 2, all non-MFJ → 1 for TAXSIM'],
        ['hourly_wage', 'pwages / hrs',
         'pwages = primary job only; hrs = all jobs → underestimates for multi-job workers',
         'Hours quality filter (520–4,160 hrs/yr) applied; caveat noted in results'],
        ['broad_income', 'Sum of 9 income components',
         'Missing NLSY codes (−1 to −5) must not be summed as negative income',
         'All negative NLSY codes → Stata missing before summing'],
    ],
    MX, by, CW, Inches(4.40),
    col_widths=[Inches(1.6), Inches(2.1), Inches(4.0), Inches(4.19)],
    row_fs=13)

slide_num(s, 14)

# ── Save ────────────────────────────────────────────────────────────────────────
out = r"D:\Stata Data\labor_signaling_project\lab_meeting_presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
